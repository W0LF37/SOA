"""Verify the v4 build against the approved v2 text."""
from pptx import Presentation
from pptx.oxml.ns import qn

prs = Presentation("output/AI_Project_Manager_Defense_v4.pptx")
print(f"Total slides: {len(prs.slides)}")
print()

total_anims = 0
for i, slide in enumerate(prs.slides, 1):
    n_shapes = len(slide.shapes)
    n_charts = sum(1 for s in slide.shapes if s.has_chart)
    n_tables = sum(1 for s in slide.shapes if s.has_table)
    n_pics = sum(1 for s in slide.shapes if s.shape_type == 13)
    has_trans = slide.element.find(qn("p:transition")) is not None
    timing = slide.element.find(qn("p:timing"))
    n_anim = 0
    if timing is not None:
        n_anim = len(timing.findall(".//" + qn("p:spTgt"))) // 2
    notes_len = len(slide.notes_slide.notes_text_frame.text) if slide.has_notes_slide else 0
    total_anims += n_anim
    print(
        f"  Slide {i:2d}: shapes={n_shapes:3d}  charts={n_charts}  tables={n_tables}  "
        f"pics={n_pics}  trans={has_trans}  anim={n_anim:2d}  notes={notes_len:4d}c"
    )

print()
print(f"Total animations across deck: {total_anims}")
print()

# Verbatim verification (v2 phrases that must be preserved)
checks = [
    (2,  "Software projects fail not for lack of ideas but because converting a brief into an executable"),
    (2,  "The architectural novelty is not the local LLM alone"),
    (4,  "type_reason explaining why it is functional or non-functional"),
    (4,  "complexity_reason explaining its complexity score"),
    (5,  "machine-readable link between requirement, task, estimate, and risk"),
    (5,  "without an audit record that a stakeholder can defend"),
    (6,  "every task carries its reasoning, the graph is validated, hours are calibrated"),
    (11, "marking status as todo, doing, or done, and asking the AI to explain any decision"),
    (13, "fully async, fully typed"),
    (14, "thirty-five percent RAG and sixty-five percent rule"),
    (14, "noisier retrieval matches and we trust the deterministic rule more"),
    (15, "graceful degradation, not crash"),
    (16, "across ten domains"),
    (16, "PRED-twenty-five by three percent"),
]

print("Verbatim v2 text verification:")
print("-" * 70)
all_pass = True
for slide_n, phrase in checks:
    notes = prs.slides[slide_n - 1].notes_slide.notes_text_frame.text
    ok = phrase in notes
    if not ok:
        all_pass = False
    marker = "PASS" if ok else "FAIL"
    print(f"  [{marker}]  Slide {slide_n:2d}  contains: {phrase[:55]}...")

print()
print(f"Overall verbatim check: {'ALL PASS' if all_pass else 'SOME FAILED'}")
