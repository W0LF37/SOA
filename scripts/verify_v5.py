"""Verify the v5 build (academic version, no animations).

Checks:
  1. 17 slides
  2. NO per-shape animations
  3. Slide transitions present
  4. Real PowerPoint tables on slides 10, 12, 13
  5. Real charts on slide 16
  6. Image on slide 14
  7. Speaker notes verbatim from v2
  8. On-slide text matches v2 verbatim
"""
from pptx import Presentation
from pptx.oxml.ns import qn

PATH = "output/AI_Project_Manager_Defense_v6.pptx"
prs = Presentation(PATH)
print(f"File: {PATH}")
print(f"Total slides: {len(prs.slides)}")
print()


def all_text(slide):
    """Concatenate ALL text on a slide from shapes (text frames + tables)."""
    chunks = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            chunks.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    chunks.append(cell.text_frame.text)
    return "\n".join(chunks)


total_anims = 0
for i, slide in enumerate(prs.slides, 1):
    n_shapes  = len(slide.shapes)
    n_charts  = sum(1 for s in slide.shapes if s.has_chart)
    n_tables  = sum(1 for s in slide.shapes if s.has_table)
    n_pics    = sum(1 for s in slide.shapes if s.shape_type == 13)
    has_trans = slide.element.find(qn("p:transition")) is not None
    timing    = slide.element.find(qn("p:timing"))
    n_anim    = 0
    if timing is not None:
        n_anim = len(timing.findall(".//" + qn("p:spTgt"))) // 2
    total_anims += n_anim
    notes_len = len(slide.notes_slide.notes_text_frame.text) if slide.has_notes_slide else 0
    print(
        f"  Slide {i:2d}: shapes={n_shapes:3d}  charts={n_charts}  tables={n_tables}  "
        f"pics={n_pics}  trans={has_trans}  anim={n_anim:2d}  notes={notes_len:4d}c"
    )

print()
print(f"Total per-shape animations: {total_anims}  (expected: 0)")
print()

# Verbatim checks — both notes AND on-slide text
notes_checks = [
    (2,  "Software projects fail not for lack of ideas but because converting a brief into an executable"),
    (2,  "The architectural novelty is not the local LLM alone"),
    (4,  "type_reason explaining why it is functional or non-functional"),
    (5,  "machine-readable link between requirement, task, estimate, and risk"),
    (5,  "without an audit record that a stakeholder can defend"),
    (14, "thirty-five percent RAG and sixty-five percent rule"),
    (16, "across ten domains"),
    (16, "PRED-twenty-five by three percent"),
]

onslide_checks = [
    (2,  "Free-text brief"),
    (2,  "Validated tasks"),
    (4,  "CORE IDEA"),
    (4,  "SEVEN CAPABILITIES"),
    (4,  "FR / NFR classification with explicit reasoning"),
    (4,  "RAG-calibrated hours"),
    (5,  "Knowledge fragmentation"),
    (5,  "Subjective decomposition"),
    (5,  "Cloud lock-in"),
    (5,  "Black-box outputs"),
    (6,  "FOUR STRATEGIC OBJECTIVES"),
    (7,  "Multi-Agent Orchestration"),
    (7,  "Hybrid Rule + LLM + RAG"),
    (8,  "THE SINGLE-PROMPT FAILURE"),
    (8,  "Planner"),
    (9,  "Gen 1"),
    (9,  "COCOMO"),
    (10, "COMPARATIVE LANDSCAPE"),
    (10, "Linear Copilot"),
    (10, "OUR CONTRIBUTION"),
    (11, "TWO ROLES"),
    (11, "Zustand"),
    (12, "FUNCTIONAL REQUIREMENTS"),
    (12, "NON-FUNCTIONAL REQUIREMENTS"),
    (12, "brief_parser.py"),
    (13, "THREE-TIER ARCHITECTURE"),
    (13, "Ollama + Qwen"),
    (14, "SIX-STAGE PIPELINE"),
    (14, "35 % RAG / 65 % rule"),
    (14, "15 % RAG / 85 % rule"),
    (15, "THREE COOPERATING LAYERS"),
    (15, "Rule Layer"),
    (15, "Desharnais"),
    (16, "ARTIFACT A"),
    (16, "ARTIFACT B"),
    (16, "ABLATION"),
    (16, "0.787"),
    (16, "100 %"),
    (17, "RESEARCH CONTRIBUTIONS"),
    (17, "ACKNOWLEDGED LIMITATIONS"),
    (17, "Thank you"),
]

print("=" * 70)
print("Speaker-notes verbatim checks:")
print("-" * 70)
notes_pass = 0
for n, phrase in notes_checks:
    txt = prs.slides[n - 1].notes_slide.notes_text_frame.text
    ok = phrase in txt
    notes_pass += ok
    marker = "PASS" if ok else "FAIL"
    print(f"  [{marker}]  Slide {n:2d}  notes contain: {phrase[:55]}")

print()
print("On-slide text checks (v2 verbatim phrases must appear on the slide):")
print("-" * 70)
slide_pass = 0
for n, phrase in onslide_checks:
    txt = all_text(prs.slides[n - 1])
    ok = phrase in txt
    slide_pass += ok
    marker = "PASS" if ok else "FAIL"
    print(f"  [{marker}]  Slide {n:2d}  contains: {phrase}")

print()
print("=" * 70)
print(f"Speaker-notes verbatim: {notes_pass}/{len(notes_checks)} passed")
print(f"On-slide verbatim:      {slide_pass}/{len(onslide_checks)} passed")
print(f"Per-shape animations:   {total_anims}  (target: 0)")
