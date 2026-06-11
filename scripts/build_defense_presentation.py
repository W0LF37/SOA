"""AI Project Manager — Graduation Defense (v6 · academic-rich, no animations).

v6 = v4 card-rich visual design  +  v2 verbatim text  +  NO animations.

What's preserved from v4:
  - Light modern theme (slate-50 background, white cards, sky+amber accents)
  - Numbered pillar cards for capabilities, pillars, agents, objectives
  - Generation-badge cards with period + name + tech + summary
  - Real PowerPoint tables (slides 10, 12, 13)
  - 2 real charts (slide 16)
  - KPI tiles (slide 16)
  - Image embed (slide 14)
  - Section badges + title accent underline (sky + amber)
  - Corner accents (subtle)

What's restored from v2:
  - All card body text uses v2's full phrasing (not shortened labels)
  - All speaker scripts verbatim v2

What's removed (vs v4):
  - All per-shape cascade animations
  - Only subtle fade slide transitions remain (academically acceptable)
"""
from __future__ import annotations

import json
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ──────────────────────────────────────────────────────────────────────────────
# THEME
# ──────────────────────────────────────────────────────────────────────────────
BG_LIGHT    = RGBColor(0xF1, 0xF5, 0xF9)
BG_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG_CARD     = RGBColor(0xFF, 0xFF, 0xFF)
INK_DEEP    = RGBColor(0x0F, 0x17, 0x2A)
INK_BODY    = RGBColor(0x33, 0x41, 0x55)
INK_MUTED   = RGBColor(0x64, 0x74, 0x8B)
INK_FAINT   = RGBColor(0x94, 0xA3, 0xB8)
BORDER      = RGBColor(0xE2, 0xE8, 0xF0)

SKY         = RGBColor(0x0E, 0xA5, 0xE9)
SKY_DARK    = RGBColor(0x02, 0x49, 0x6B)
SKY_LIGHT   = RGBColor(0xE0, 0xF2, 0xFE)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
EMERALD     = RGBColor(0x10, 0xB9, 0x81)
ROSE        = RGBColor(0xF4, 0x3F, 0x5E)
VIOLET      = RGBColor(0x8B, 0x5C, 0xF6)

FONT_TITLE = "Segoe UI"
FONT_BODY  = "Segoe UI"
FONT_MONO  = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

REPO         = Path(__file__).resolve().parents[1]
OUT_PATH     = REPO / "output" / "AI_Project_Manager_Defense_v6.pptx"
PIPELINE_IMG = REPO / "poster_assets" / "pipeline_flow.png"


# ──────────────────────────────────────────────────────────────────────────────
# Generic helpers
# ──────────────────────────────────────────────────────────────────────────────
def set_background(slide, color=BG_LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_corner_accent(slide):
    tri = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, Inches(11.8), Inches(-0.4), Inches(1.7), Inches(1.7))
    tri.line.fill.background()
    tri.fill.solid()
    tri.fill.fore_color.rgb = SKY
    tri.rotation = 90.0
    tri2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, Inches(12.7), Inches(7.0), Inches(0.9), Inches(0.9))
    tri2.line.fill.background()
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = AMBER
    tri2.rotation = 180.0


def add_section_badge(slide, section_text):
    if not section_text:
        return
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.35), Inches(1.7), Inches(0.42))
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = SKY_LIGHT
    badge.line.color.rgb = SKY
    badge.line.width = Pt(0.75)
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = section_text
    r.font.name = FONT_BODY
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = SKY_DARK


def add_title(slide, text, subtitle=None, section=None):
    if section:
        add_section_badge(slide, section)
        title_top = Inches(0.85)
    else:
        title_top = Inches(0.45)
    tb = slide.shapes.add_textbox(Inches(0.6), title_top, Inches(12.0), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT_TITLE
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = INK_DEEP
    bar1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), title_top + Inches(0.62), Inches(0.85), Inches(0.06))
    bar1.line.fill.background()
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = SKY
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), title_top + Inches(0.62), Inches(0.35), Inches(0.06))
    bar2.line.fill.background()
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = AMBER
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.6), title_top + Inches(0.72), Inches(12.0), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = FONT_BODY
        sr.font.size = Pt(13)
        sr.font.color.rgb = INK_MUTED
        sr.font.italic = True


def add_footer(slide, n, total=17):
    fb = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(11.0), Inches(0.3))
    p = fb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "AI Project Manager  ·  Graduation Defense  ·  2026"
    r.font.name = FONT_BODY
    r.font.size = Pt(9)
    r.font.color.rgb = INK_FAINT
    pb = slide.shapes.add_textbox(Inches(11.7), Inches(7.1), Inches(1.1), Inches(0.3))
    pp = pb.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    pr = pp.add_run()
    pr.text = f"{n} / {total}"
    pr.font.name = FONT_BODY
    pr.font.size = Pt(10)
    pr.font.bold = True
    pr.font.color.rgb = SKY


def add_card(slide, left, top, width, height, fill=BG_CARD, border=BORDER, border_w=0.75, radius=0.04):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = radius
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = border
        card.line.width = Pt(border_w)
    card.shadow.inherit = False
    return card


def add_text_in_card(card, paragraphs, padding_left=Inches(0.25), padding_top=Inches(0.2)):
    tf = card.text_frame
    tf.margin_left = padding_left
    tf.margin_right = padding_left
    tf.margin_top = padding_top
    tf.margin_bottom = padding_top
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (text, size, bold, color, align, indent) in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if indent:
            p.level = indent
        r = p.add_run()
        r.text = text
        r.font.name = FONT_BODY
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def add_panel_title(slide, left, top, text, color=INK_DEEP, size=14, width=None):
    w = width if width else Inches(4.0)
    tb = slide.shapes.add_textbox(left, top, w, Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT_TITLE
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color


def add_kpi_tile(slide, left, top, width, height, big, label, accent=SKY, bg=BG_WHITE):
    card = add_card(slide, left, top, width, height, fill=bg, border=BORDER)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    tf = card.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = big
    r1.font.name = FONT_TITLE
    r1.font.size = Pt(30)
    r1.font.bold = True
    r1.font.color.rgb = INK_DEEP
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = FONT_BODY
    r2.font.size = Pt(10)
    r2.font.color.rgb = INK_MUTED
    return card


def add_numbered_card(slide, left, top, width, height, num, title, body, accent, body_size=11):
    card = add_card(slide, left, top, width, height, fill=BG_WHITE, border=BORDER)
    # numbered circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.2), Inches(0.45), Inches(0.45))
    circle.line.fill.background()
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.shadow.inherit = False
    ct = circle.text_frame
    ct.margin_left = ct.margin_right = 0
    ct.margin_top = ct.margin_bottom = 0
    ct.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ct.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = str(num)
    cr.font.name = FONT_TITLE
    cr.font.size = Pt(16)
    cr.font.bold = True
    cr.font.color.rgb = BG_WHITE
    # title
    tb = slide.shapes.add_textbox(
        left + Inches(0.75), top + Inches(0.2), width - Inches(0.95), Inches(0.5))
    tp = tb.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.name = FONT_TITLE
    tr.font.size = Pt(14)
    tr.font.bold = True
    tr.font.color.rgb = INK_DEEP
    # body
    bb = slide.shapes.add_textbox(
        left + Inches(0.25), top + Inches(0.78),
        width - Inches(0.5), height - Inches(0.95))
    bf = bb.text_frame
    bf.word_wrap = True
    bp = bf.paragraphs[0]
    bp.line_spacing = 1.2
    br = bp.add_run()
    br.text = body
    br.font.name = FONT_BODY
    br.font.size = Pt(body_size)
    br.font.color.rgb = INK_BODY
    return card


# ──────────────────────────────────────────────────────────────────────────────
# Charts
# ──────────────────────────────────────────────────────────────────────────────
def style_chart_series(chart, colors, label_color=INK_BODY):
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(9)
    dl.font.color.rgb = label_color
    dl.font.name = FONT_BODY
    for i, series in enumerate(plot.series):
        f = series.format.fill
        f.solid()
        f.fore_color.rgb = colors[i % len(colors)]
        series.format.line.fill.background()
    for axis in (chart.category_axis, chart.value_axis):
        axis.tick_labels.font.size = Pt(9)
        axis.tick_labels.font.color.rgb = INK_MUTED
        axis.tick_labels.font.name = FONT_BODY
    if chart.has_legend:
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = INK_BODY
        chart.legend.font.name = FONT_BODY
        chart.legend.include_in_layout = False
    if chart.has_title:
        tr = chart.chart_title.text_frame.paragraphs[0].runs[0]
        tr.font.color.rgb = INK_DEEP
        tr.font.size = Pt(13)
        tr.font.bold = True
        tr.font.name = FONT_TITLE


def add_threshold_chart(slide, left, top, width, height):
    cd = CategoryChartData()
    cd.categories = ["Pass Rate", "Critic Score", "Overall Score", "Direct LLM"]
    cd.add_series("Achieved",  (1.00, 0.986, 0.787, 0.75))
    cd.add_series("Threshold", (1.00, 0.85,  0.78,  0.70))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd)
    c = gf.chart
    c.has_title = True
    c.chart_title.text_frame.text = "All 4 Quality Thresholds Met"
    c.has_legend = True
    c.legend.position = XL_LEGEND_POSITION.BOTTOM
    style_chart_series(c, [SKY, AMBER])
    return gf


def add_ablation_chart(slide, left, top, width, height):
    cd = CategoryChartData()
    cd.categories = ["MMRE (lower = better)", "PRED(25) (higher = better)"]
    cd.add_series("R: Rules only", (0.271, 0.778))
    cd.add_series("K: Rules + KB", (0.255, 0.755))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd)
    c = gf.chart
    c.has_title = True
    c.chart_title.text_frame.text = "Ablation: KB Impact on Effort Estimation"
    c.has_legend = True
    c.legend.position = XL_LEGEND_POSITION.BOTTOM
    style_chart_series(c, [INK_FAINT, SKY])
    return gf


# ──────────────────────────────────────────────────────────────────────────────
# Notes + (subtle) slide transitions only — NO per-shape animations
# ──────────────────────────────────────────────────────────────────────────────
def add_notes(slide, script):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = script
    r.font.name = FONT_BODY
    r.font.size = Pt(14)


def add_slide_transition(slide, kind="fade"):
    sld = slide.element
    cSld = sld.find(qn("p:cSld"))
    if cSld is None:
        return
    existing = sld.find(qn("p:transition"))
    if existing is not None:
        sld.remove(existing)
    trans = etree.SubElement(sld, qn("p:transition"))
    trans.set("spd", "med")
    trans.set("advClick", "1")
    if kind == "fade":
        etree.SubElement(trans, qn("p:fade"))
    sld.remove(trans)
    cSld.addnext(trans)


# ──────────────────────────────────────────────────────────────────────────────
# Slide builders (v2 verbatim text inside v4-style cards)
# ──────────────────────────────────────────────────────────────────────────────
def build_cover(slide):
    set_background(slide, BG_WHITE)
    add_card(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT, border=None, radius=0)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = SKY
    stripe2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.32), Inches(8), Inches(0.18))
    stripe2.line.fill.background()
    stripe2.fill.solid()
    stripe2.fill.fore_color.rgb = AMBER

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "AI PROJECT MANAGER"
    r.font.name = FONT_TITLE
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.color.rgb = INK_DEEP

    sb = slide.shapes.add_textbox(Inches(0.85), Inches(3.55), Inches(11.5), Inches(0.55))
    sp = sb.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = "A Local Multi-Agent System for"
    sr.font.name = FONT_BODY
    sr.font.size = Pt(22)
    sr.font.color.rgb = INK_MUTED

    sb2 = slide.shapes.add_textbox(Inches(0.85), Inches(4.0), Inches(11.5), Inches(0.55))
    sp = sb2.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = "Automated Software Project Planning"
    sr.font.name = FONT_BODY
    sr.font.size = Pt(22)
    sr.font.color.rgb = INK_MUTED

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(4.85), Inches(2.2), Inches(0.08))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = SKY

    fb = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.5))
    ff = fb.text_frame
    ff.word_wrap = True
    fp = ff.paragraphs[0]
    fr = fp.add_run()
    fr.text = "[Student Name]    ·    [Supervisor Name]"
    fr.font.name = FONT_BODY
    fr.font.size = Pt(16)
    fr.font.color.rgb = INK_BODY
    fp2 = ff.add_paragraph()
    fr2 = fp2.add_run()
    fr2.text = "Department of Informatics Engineering    ·    Graduation Defense 2026"
    fr2.font.name = FONT_BODY
    fr2.font.size = Pt(14)
    fr2.font.bold = True
    fr2.font.color.rgb = SKY


def build_intro(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Project Introduction",
              "From a free-text brief to a complete, auditable plan — locally")

    # Tagline card
    intro = add_card(slide, Inches(0.6), Inches(1.85), Inches(12.1), Inches(1.5),
                     fill=SKY_LIGHT, border=SKY, border_w=1.0)
    add_text_in_card(intro, [
        ("We built a coordinated team of four AI agents that converts a free-text project brief",
         14, False, INK_DEEP, PP_ALIGN.CENTER, 0),
        ("into a complete, auditable execution plan — all running locally via Ollama.",
         14, True, INK_DEEP, PP_ALIGN.CENTER, 0),
    ], padding_top=Inches(0.3))

    # Input / Output cards
    add_panel_title(slide, Inches(0.6), Inches(3.55), "INPUT", SKY, size=12)
    in_card = add_card(slide, Inches(0.6), Inches(3.9), Inches(3.3), Inches(2.85))
    add_text_in_card(in_card, [
        ("📄  Free-text brief", 18, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("    ", 14, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("Any plain-language requirements", 11, False, INK_MUTED, PP_ALIGN.LEFT, 0),
        ("document — clinic, fintech, LMS,", 11, False, INK_MUTED, PP_ALIGN.LEFT, 0),
        ("hospital, e-commerce, …", 11, False, INK_MUTED, PP_ALIGN.LEFT, 0),
    ])

    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(4.05), Inches(5.05), Inches(0.6), Inches(0.5))
    arrow.line.fill.background()
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = SKY

    add_panel_title(slide, Inches(4.85), Inches(3.55), "OUTPUT — FULL EXECUTION PLAN", SKY, size=12)
    out_card = add_card(slide, Inches(4.85), Inches(3.9), Inches(7.85), Inches(2.85))
    add_text_in_card(out_card, [
        ("✓  Validated tasks  (with reasoning)",                12, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("✓  Dependency DAG with critical path",                12, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("✓  Calibrated effort estimates",                      12, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("✓  Risk report  (six categories)",                    12, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("✓  Sprint schedule",                                  12, False, INK_BODY, PP_ALIGN.LEFT, 0),
    ], padding_top=Inches(0.2))

    add_footer(slide, idx)


def build_toc(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Table of Contents")
    sections = [
        ("1", "Introduction",                "Idea · Problem · Objectives",      SKY),
        ("2", "Literature & Related Work",   "Evolution · Research gap",         AMBER),
        ("3", "System Analysis",             "Actors · Requirements · Stack",    EMERALD),
        ("4", "AI Design",                   "Multi-agent · Hybrid reasoning",   VIOLET),
        ("5", "Evaluation",                  "20 samples · Ablation study",      ROSE),
    ]
    top = Inches(1.85)
    h = Inches(0.85)
    for i, (n, title, sub, color) in enumerate(sections):
        y = top + i * (h + Inches(0.12))
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.15), Inches(0.55), Inches(0.55))
        circle.line.fill.background()
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        ct = circle.text_frame
        ct.margin_left = ct.margin_right = 0
        ct.margin_top = ct.margin_bottom = 0
        ct.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ct.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = n
        cr.font.name = FONT_TITLE
        cr.font.size = Pt(18)
        cr.font.bold = True
        cr.font.color.rgb = BG_WHITE
        tb = slide.shapes.add_textbox(Inches(1.6), y + Inches(0.07), Inches(11.0), Inches(0.45))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = FONT_TITLE
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = INK_DEEP
        sb = slide.shapes.add_textbox(Inches(1.6), y + Inches(0.45), Inches(11.0), Inches(0.35))
        sp = sb.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = sub
        sr.font.name = FONT_BODY
        sr.font.size = Pt(12)
        sr.font.color.rgb = INK_MUTED
    add_footer(slide, idx)


def build_idea(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "The Core Idea & Capabilities", section="§1  Introduction")

    # CORE IDEA panel
    add_panel_title(slide, Inches(0.6), Inches(1.6), "CORE IDEA", SKY, size=12)
    idea = add_card(slide, Inches(0.6), Inches(1.95), Inches(12.1), Inches(1.0),
                    fill=SKY_LIGHT, border=SKY, border_w=1.0)
    add_text_in_card(idea, [
        ("A planner that thinks — not a template that fills.",
         17, True, INK_DEEP, PP_ALIGN.CENTER, 0),
        ("Four agents that decompose, validate, schedule, and monitor a software project end-to-end, locally.",
         12, False, INK_BODY, PP_ALIGN.CENTER, 0),
    ], padding_top=Inches(0.15))

    add_panel_title(slide, Inches(0.6), Inches(3.05), "SEVEN CAPABILITIES", SKY, size=12)
    # 7 capability cards — v2 verbatim phrases
    caps = [
        ("Parse",       "FR / NFR classification with explicit reasoning",       SKY),
        ("Decompose",   "Rule-based + LLM-assisted task breakdown",              SKY),
        ("Estimate",    "RAG-calibrated hours  (rule-dominant blend)",           SKY),
        ("Sequence",    "NetworkX DAG  ·  critical path  ·  parallel groups",    SKY),
        ("Risk-Score",  "Six categories, rule layer + LLM novel risks",          AMBER),
        ("Schedule",    "Topological sprints with team capacity",                AMBER),
        ("Monitor",     "Git commit ↔ task semantic matching",                   AMBER),
    ]
    top1 = Inches(3.4)
    top2 = Inches(5.25)
    tile_w = Inches(2.93)
    tile_h = Inches(1.7)
    for i in range(4):
        left = Inches(0.6 + i * 3.06)
        h, b, c = caps[i]
        add_numbered_card(slide, left, top1, tile_w, tile_h, i+1, h, b, c, body_size=10)
    for j in range(3):
        i = 4 + j
        left = Inches(2.13 + j * 3.06)
        h, b, c = caps[i]
        add_numbered_card(slide, left, top2, tile_w, tile_h, i+1, h, b, c, body_size=10)
    add_footer(slide, idx)


def build_problem(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Problem Statement", section="§1  Introduction")

    tag = add_card(slide, Inches(0.6), Inches(1.85), Inches(12.1), Inches(0.95),
                   fill=AMBER_LIGHT, border=AMBER, border_w=1.0)
    add_text_in_card(tag, [
        ("Manual project planning is slow, inconsistent, and gated behind senior expertise.",
         13, False, INK_DEEP, PP_ALIGN.CENTER, 0),
        ("Current AI tools sacrifice privacy and auditability for convenience.",
         13, True, INK_DEEP, PP_ALIGN.CENTER, 0),
    ], padding_top=Inches(0.15))

    causes = [
        ("Knowledge fragmentation",
         "Briefs, rules, and historical data live in disconnected files.",
         SKY,    "📁"),
        ("Subjective decomposition",
         "Two planners produce two different plans from the same brief.",
         AMBER,  "🔀"),
        ("Cloud lock-in & privacy risk",
         "Jira AI · Linear Copilot · GitHub Workspace — all cloud-only.",
         ROSE,   "🌐"),
        ("Black-box outputs",
         "No reasoning trail, no fallback path, no audit record.",
         VIOLET, "⬛"),
    ]
    cw = Inches(5.95)
    ch = Inches(1.85)
    positions = [(0.6, 3.05), (6.75, 3.05), (0.6, 5.05), (6.75, 5.05)]
    for (cx, cy), (head, body, color, icon) in zip(positions, causes):
        card = add_card(slide, Inches(cx), Inches(cy), cw, ch, fill=BG_WHITE, border=BORDER)
        # icon circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx + 0.25), Inches(cy + 0.25), Inches(0.55), Inches(0.55))
        circle.line.fill.background()
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        ib = slide.shapes.add_textbox(Inches(cx + 0.25), Inches(cy + 0.25), Inches(0.55), Inches(0.55))
        itf = ib.text_frame
        itf.margin_left = itf.margin_right = 0
        itf.margin_top = itf.margin_bottom = 0
        itf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ip = itf.paragraphs[0]
        ip.alignment = PP_ALIGN.CENTER
        ir = ip.add_run()
        ir.text = icon
        ir.font.name = "Segoe UI Emoji"
        ir.font.size = Pt(18)
        # head
        hb = slide.shapes.add_textbox(Inches(cx + 0.95), Inches(cy + 0.25), cw - Inches(1.0), Inches(0.45))
        hp = hb.text_frame.paragraphs[0]
        hr = hp.add_run()
        hr.text = head
        hr.font.name = FONT_TITLE
        hr.font.size = Pt(14)
        hr.font.bold = True
        hr.font.color.rgb = INK_DEEP
        # body
        bb = slide.shapes.add_textbox(Inches(cx + 0.25), Inches(cy + 0.95), cw - Inches(0.5), Inches(0.9))
        bf = bb.text_frame
        bf.word_wrap = True
        bp = bf.paragraphs[0]
        bp.line_spacing = 1.2
        br = bp.add_run()
        br.text = body
        br.font.name = FONT_BODY
        br.font.size = Pt(11)
        br.font.color.rgb = INK_BODY
    add_footer(slide, idx)


def build_objectives(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Project Objectives", section="§1  Introduction")
    add_panel_title(slide, Inches(0.6), Inches(1.6), "FOUR STRATEGIC OBJECTIVES", SKY, size=12)
    objs = [
        ("Build a comprehensive auditable planner",
         "Tasks with reasoning fields · validated DAG · calibrated hours · categorized risk — from one brief.",
         SKY),
        ("Compare hybrid AI architectures empirically",
         "Quantify rule vs LLM vs RAG contributions on 20 ground-truth samples  (ablation methodology).",
         AMBER),
        ("Guarantee reliability through a deterministic safety layer",
         "Critic Agent · independent validation · rejection threshold at score 0.50.",
         EMERALD),
        ("Optimize local-performance / quality trade-off",
         "No paid APIs  ·  MMRE ≈ 0.36  ·  PRED(25) ≈ 0.67  on a single workstation.",
         VIOLET),
    ]
    top = Inches(2.0)
    h = Inches(1.15)
    for i, (head, body, color) in enumerate(objs):
        y = top + i * (h + Inches(0.13))
        add_numbered_card(slide, Inches(0.6), y, Inches(12.1), h, i+1, head, body, color, body_size=12)
    add_footer(slide, idx)


def build_pillars(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Proposed System — Five Pillars", section="§1  Introduction")
    pillars = [
        ("Multi-Agent Orchestration",
         "Four agents with narrow contracts — not one monolithic prompt.",
         SKY),
        ("Local-First Operation",
         "Qwen 9B Q5_K_M via Ollama  ·  zero paid API calls once weights and embeddings are cached locally.",
         AMBER),
        ("Hybrid Rule + LLM + RAG Reasoning",
         "Deterministic rules  ·  LLM Planning Reasoner  ·  ChromaDB knowledge base (100 curated documents).",
         EMERALD),
        ("Dependency-Aware Output",
         "NetworkX DAG  ·  NFR-to-FR constraint propagation  ·  critical path  ·  bottleneck analysis.",
         VIOLET),
        ("Auditable Evaluation Methodology",
         "20 ground-truth samples  ·  F1-FR · F1-NFR · MMRE · PRED(25)  ·  2-condition ablation (Rules vs Rules + KB).",
         ROSE),
    ]
    cw = Inches(3.93)
    ch = Inches(2.3)
    top_row_y = Inches(1.85)
    bot_row_y = Inches(4.45)
    for i in range(3):
        left = Inches(0.6 + i * 4.07)
        h, b, c = pillars[i]
        add_numbered_card(slide, left, top_row_y, cw, ch, i+1, h, b, c, body_size=11)
    for j in range(2):
        i = 3 + j
        left = Inches(2.63 + j * 4.07)
        h, b, c = pillars[i]
        add_numbered_card(slide, left, bot_row_y, cw, ch, i+1, h, b, c, body_size=11)
    add_footer(slide, idx)


def build_multiagent(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Why Multi-Agent (Not a Single Prompt)", section="§1  Introduction")

    warn = add_card(slide, Inches(0.6), Inches(1.85), Inches(12.1), Inches(1.0),
                    fill=AMBER_LIGHT, border=AMBER, border_w=1.0)
    add_text_in_card(warn, [
        ("THE SINGLE-PROMPT FAILURE", 12, True, AMBER, PP_ALIGN.CENTER, 0),
        ("A single LLM call cannot simultaneously decompose, validate, propagate constraints,",
         12, False, INK_BODY, PP_ALIGN.CENTER, 0),
        ("and score risk without hallucinating in at least one dimension.",
         12, False, INK_BODY, PP_ALIGN.CENTER, 0),
    ], padding_top=Inches(0.1))

    agents = [
        ("Planner",
         "Decomposes requirements into tasks · assigns type_reason · complexity_reason · tags semantically (auth, crud, security, …).",
         SKY),
        ("Critic",
         "Independent validator (rules + LLM layers) · penalties: error 0.20 · warning 0.05 · info 0.01 · halts pipeline at score < 0.50.",
         ROSE),
        ("Risk",
         "Six categories: bottleneck · complexity · schedule · dependency · resource · quality.",
         AMBER),
        ("Monitor",
         "Git commit ↔ task semantic matching (MiniLM-L6-v2 · threshold 0.65) + keyword fallback.",
         EMERALD),
    ]
    cw = Inches(2.93)
    ch = Inches(3.6)
    for i, (h, b, c) in enumerate(agents):
        left = Inches(0.6 + i * 3.06)
        add_numbered_card(slide, left, Inches(3.05), cw, ch, i+1, h, b, c, body_size=10)
    add_footer(slide, idx)


def build_generations(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Three Generations of Planning AI",
              section="§2  Literature & Related Work")
    gens = [
        ("Gen 1", "1981 — 2000", "Algorithmic Estimation",
         "COCOMO · COCOMO II · Function Points",
         "✓  Mathematically defensible       ✗  No task generation · no risk analysis · expert-only",
         INK_FAINT),
        ("Gen 2", "2000 — 2020", "Classical Machine Learning",
         "Regression effort estimation · SVM risk classifiers",
         "✓  Data-driven calibration         ✗  No decomposition · no reasoning · black-box",
         AMBER),
        ("Gen 3", "2023 — Now", "LLM & RAG-Based Planning",
         "GPT-4 plan generation · Jira AI · Linear Copilot · GitHub Workspace · AutoGPT-style agents",
         "✓  Natural-language input · generates tasks    ✗  Hallucination · privacy · no auditable reasoning · no deterministic safety layer",
         SKY),
    ]
    h = Inches(1.55)
    for i, (gen, period, name, tech, summary, color) in enumerate(gens):
        y = Inches(1.85) + i * (h + Inches(0.18))
        add_card(slide, Inches(0.6), y, Inches(12.1), h)
        # badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), y + Inches(0.2),
            Inches(1.4), Inches(1.15))
        badge.adjustments[0] = 0.15
        badge.line.fill.background()
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        bt = badge.text_frame
        bt.margin_left = bt.margin_right = 0
        bt.margin_top = bt.margin_bottom = 0
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp1 = bt.paragraphs[0]
        bp1.alignment = PP_ALIGN.CENTER
        br1 = bp1.add_run()
        br1.text = gen
        br1.font.name = FONT_TITLE
        br1.font.size = Pt(20)
        br1.font.bold = True
        br1.font.color.rgb = BG_WHITE
        bp2 = bt.add_paragraph()
        bp2.alignment = PP_ALIGN.CENTER
        br2 = bp2.add_run()
        br2.text = period
        br2.font.name = FONT_BODY
        br2.font.size = Pt(10)
        br2.font.color.rgb = BG_WHITE
        # right text
        tb = slide.shapes.add_textbox(Inches(2.5), y + Inches(0.2), Inches(10.0), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = name
        r1.font.name = FONT_TITLE
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = INK_DEEP
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = tech
        r2.font.name = FONT_BODY
        r2.font.size = Pt(12)
        r2.font.italic = True
        r2.font.color.rgb = INK_MUTED
        p3 = tf.add_paragraph()
        p3.line_spacing = 1.2
        r3 = p3.add_run()
        r3.text = summary
        r3.font.name = FONT_BODY
        r3.font.size = Pt(11)
        r3.font.color.rgb = INK_BODY
    add_footer(slide, idx)


def _style_table_header(cell, text):
    cell.fill.solid()
    cell.fill.fore_color.rgb = SKY_DARK
    tf = cell.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = BG_WHITE


def _style_table_body(cell, text, *, zebra=False, mono=False, bold=False, size=10):
    cell.fill.solid()
    cell.fill.fore_color.rgb = BG_LIGHT if zebra else BG_WHITE
    tf = cell.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT_MONO if mono else FONT_BODY
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = INK_BODY


def build_relwork(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Related Work & Research Gap", section="§2  Literature & Related Work")

    add_panel_title(slide, Inches(0.6), Inches(1.7), "COMPARATIVE LANDSCAPE", SKY)
    tbl_shape = slide.shapes.add_table(7, 4, Inches(0.6), Inches(2.05), Inches(7.7), Inches(2.85))
    table = tbl_shape.table
    headers = ["Tool / Reference", "Year", "Approach", "Key Gap We Address"]
    data = [
        ("COCOMO II",            "2000", "Algorithmic",   "No task generation"),
        ("PROMISE NASA dataset", "2005", "ML estimation", "No decomposition"),
        ("Jira AI Backlog Gen.", "2023", "Cloud LLM",     "Privacy · no audit"),
        ("Linear Copilot",       "2024", "Cloud LLM",     "Privacy · no rules"),
        ("GitHub Copilot WS",    "2024", "Cloud agent",   "Privacy · no NFR prop."),
        ("AutoGPT-style agents", "2024", "Recursive LLM", "No determinism"),
    ]
    widths = [Inches(2.4), Inches(0.8), Inches(1.6), Inches(2.9)]
    for j, w in enumerate(widths):
        table.columns[j].width = w
    for j, h in enumerate(headers):
        _style_table_header(table.cell(0, j), h)
    for i, row in enumerate(data, start=1):
        zebra = (i % 2 == 0)
        for j, val in enumerate(row):
            _style_table_body(table.cell(i, j), val, zebra=zebra, size=10)

    add_panel_title(slide, Inches(8.55), Inches(1.7), "OUR CONTRIBUTION  ·  4 GAPS CLOSED", AMBER)
    contrib = add_card(slide, Inches(8.55), Inches(2.05), Inches(4.15), Inches(4.6),
                       fill=BG_WHITE, border=AMBER, border_w=1.0)
    add_text_in_card(contrib, [
        ("Locality",      12, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("A fully-local multi-agent planner prototype", 10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 4, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("Auditability",  12, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("Explicit type_reason · complexity_reason on every task",
         10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 4, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("Safety",        12, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("Independent Critic with hard rejection authority",
         10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 4, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("NFR Semantics", 12, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("Automated propagation of security · offline · compliance NFRs as constraints on downstream FRs",
         10, False, INK_BODY, PP_ALIGN.LEFT, 0),
    ])
    add_footer(slide, idx)


def build_actors(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "System Actors & Use-Case Boundaries", section="§3  System Analysis")
    add_panel_title(slide, Inches(0.6), Inches(1.6), "TWO ROLES — UI-LEVEL SEPARATION (PROTOTYPE)", SKY, size=12)

    add_card(slide, Inches(0.6), Inches(1.95), Inches(5.95), Inches(4.25))
    add_card(slide, Inches(6.75), Inches(1.95), Inches(5.95), Inches(4.25))

    # Student header
    s_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.95), Inches(5.95), Inches(0.5))
    s_strip.line.fill.background()
    s_strip.fill.solid()
    s_strip.fill.fore_color.rgb = SKY
    s_tb = slide.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(5.5), Inches(0.4))
    sp = s_tb.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = "👤  STUDENT     ·     Submits a project brief and tracks plan execution"
    sr.font.name = "Segoe UI Emoji"
    sr.font.size = Pt(12)
    sr.font.bold = True
    sr.font.color.rgb = BG_WHITE

    student_items = [
        "Submit free-text brief",
        "Track task progress (Todo / Doing / Done)",
        "Rate tasks 1–5 stars (feedback loop)",
        "Ask AI to explain any task",
        "View tech stack · sprints · risks",
    ]
    s_body = slide.shapes.add_textbox(Inches(0.85), Inches(2.65), Inches(5.5), Inches(3.4))
    s_tf = s_body.text_frame
    s_tf.word_wrap = True
    for i, item in enumerate(student_items):
        p = s_tf.paragraphs[0] if i == 0 else s_tf.add_paragraph()
        p.line_spacing = 1.4
        m = p.add_run()
        m.text = "▸  "
        m.font.name = FONT_BODY
        m.font.size = Pt(13)
        m.font.bold = True
        m.font.color.rgb = SKY
        r = p.add_run()
        r.text = item
        r.font.name = FONT_BODY
        r.font.size = Pt(13)
        r.font.color.rgb = INK_BODY

    # Supervisor header
    sup_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.75), Inches(1.95), Inches(5.95), Inches(0.5))
    sup_strip.line.fill.background()
    sup_strip.fill.solid()
    sup_strip.fill.fore_color.rgb = VIOLET
    sup_tb = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.4))
    sp = sup_tb.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = "🎓  SUPERVISOR  ·  Reviews, approves, and exports the plan"
    sr.font.name = "Segoe UI Emoji"
    sr.font.size = Pt(12)
    sr.font.bold = True
    sr.font.color.rgb = BG_WHITE

    sup_items = [
        "Triage Admin-Review queue",
        "Approve / Edit / Reject flagged tasks",
        "Approve or reject full plan via in-app chat",
        "Export committee brief as Excel",
        "Run evaluation suite · view ablation",
    ]
    su_body = slide.shapes.add_textbox(Inches(7.0), Inches(2.65), Inches(5.5), Inches(3.4))
    su_tf = su_body.text_frame
    su_tf.word_wrap = True
    for i, item in enumerate(sup_items):
        p = su_tf.paragraphs[0] if i == 0 else su_tf.add_paragraph()
        p.line_spacing = 1.4
        m = p.add_run()
        m.text = "▸  "
        m.font.name = FONT_BODY
        m.font.size = Pt(13)
        m.font.bold = True
        m.font.color.rgb = VIOLET
        r = p.add_run()
        r.text = item
        r.font.name = FONT_BODY
        r.font.size = Pt(13)
        r.font.color.rgb = INK_BODY

    # SCOPE NOTE
    note = add_card(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.7),
                    fill=AMBER_LIGHT, border=AMBER, border_w=0.75)
    add_text_in_card(note, [
        ("SCOPE NOTE  ·  The prototype uses a login page with hardcoded demo credentials for role separation in the UI layer (Zustand store).",
         11, False, INK_DEEP, PP_ALIGN.CENTER, 0),
        ("Production deployment would require JWT-protected API endpoints — noted as future work.",
         11, False, INK_DEEP, PP_ALIGN.CENTER, 0),
    ], padding_top=Inches(0.1))
    add_footer(slide, idx)


def build_requirements(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Functional & Non-Functional Requirements", section="§3  System Analysis")

    add_panel_title(slide, Inches(0.6), Inches(1.7), "FUNCTIONAL REQUIREMENTS", SKY)
    fr_tbl = slide.shapes.add_table(9, 3, Inches(0.6), Inches(2.05), Inches(7.4), Inches(4.5)).table
    fr_widths = [Inches(1.0), Inches(3.6), Inches(2.8)]
    for j, w in enumerate(fr_widths):
        fr_tbl.columns[j].width = w
    fr_data = [
        ("ID", "Requirement", "Module"),
        ("REQ-F1", "Free-text brief parsing", "brief_parser.py"),
        ("REQ-F2", "Multi-agent task generation", "planner.py"),
        ("REQ-F3", "Two-layer plan validation", "critic.py"),
        ("REQ-F4", "RAG-calibrated effort estimation", "effort_estimator.py"),
        ("REQ-F5", "Dependency DAG with critical path", "dependency_graph.py"),
        ("REQ-F6", "Sprint planning with team capacity", "sprint_planner.py"),
        ("REQ-F7", "Six-category risk analysis", "risk_analyzer.py"),
        ("REQ-F8", "Git-based progress monitoring", "monitor.py"),
    ]
    for j, h in enumerate(fr_data[0]):
        _style_table_header(fr_tbl.cell(0, j), h)
    for i, row in enumerate(fr_data[1:], start=1):
        zebra = (i % 2 == 0)
        for j, val in enumerate(row):
            mono = (j == 2)
            bold = (j == 0)
            _style_table_body(fr_tbl.cell(i, j), val, zebra=zebra, mono=mono, bold=bold, size=9)

    add_panel_title(slide, Inches(8.25), Inches(1.7), "NON-FUNCTIONAL REQUIREMENTS", AMBER)
    nfr_card = add_card(slide, Inches(8.25), Inches(2.05), Inches(4.45), Inches(4.5))
    nfr_items = [
        ("NFR-1   Locality",        "No external API calls during operation"),
        ("NFR-2   Determinism",     "Temperature = 0  ·  rule-based fallback"),
        ("NFR-3   Auditability",    "Reasoning fields on every output"),
        ("NFR-4   Reproducibility", "Frozen 20-sample evaluation benchmark"),
        ("NFR-5   Safety",          "Critic rejection at score < 0.50"),
        ("NFR-6   Resilience",      "OOM-aware context downscaling"),
    ]
    tf = nfr_card.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True
    for i, (head, body) in enumerate(nfr_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        r1 = p.add_run()
        r1.text = head
        r1.font.name = FONT_BODY
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = INK_DEEP
        r2 = p.add_run()
        r2.text = "\n" + body
        r2.font.name = FONT_BODY
        r2.font.size = Pt(10)
        r2.font.color.rgb = INK_MUTED
        if i < len(nfr_items) - 1:
            sp = tf.add_paragraph()
            spr = sp.add_run()
            spr.text = " "
            spr.font.size = Pt(4)
    add_footer(slide, idx)


def build_architecture(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Architecture & Technology Stack", section="§3  System Analysis")
    add_panel_title(slide, Inches(0.6), Inches(1.6), "THREE-TIER ARCHITECTURE", SKY, size=12)
    tiers = [
        ("PRESENTATION TIER",
         "React 19  ·  TypeScript  ·  Vite  ·  Zustand  ·  13 pages  ·  @xyflow/react graph  ·  SSE streaming",
         SKY),
        ("SERVICE TIER",
         "FastAPI  ·  Python 3.14  ·  10 routers  ·  30+ endpoints  ·  Swagger documented",
         EMERALD),
        ("INTELLIGENCE TIER",
         "Ollama + Qwen 9B Q5_K_M  (temp = 0, ctx = 4096)  ·  ChromaDB 1.5.8  ·  all-MiniLM-L6-v2 embeddings  ·  NetworkX  ·  Pydantic v2",
         VIOLET),
    ]
    top = Inches(2.0)
    h = Inches(0.85)
    for i, (head, body, color) in enumerate(tiers):
        y = top + i * (h + Inches(0.12))
        add_card(slide, Inches(0.6), y, Inches(12.1), h, fill=BG_WHITE, border=BORDER)
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.12), h)
        strip.line.fill.background()
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        tb = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.1), Inches(11.7), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = head
        r.font.name = FONT_TITLE
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = color
        bb = slide.shapes.add_textbox(Inches(0.85), y + Inches(0.42), Inches(11.7), Inches(0.45))
        bp = bb.text_frame.paragraphs[0]
        br = bp.add_run()
        br.text = body
        br.font.name = FONT_BODY
        br.font.size = Pt(11)
        br.font.color.rgb = INK_BODY

    add_panel_title(slide, Inches(0.6), Inches(4.95), "TECHNOLOGY DECISIONS", SKY)
    tbl = slide.shapes.add_table(6, 3, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.65)).table
    headers = ["Choice", "Rationale", "Rejected"]
    data = [
        ("Ollama + Qwen", "Local  ·  deterministic",   "GPT-4 (paid · online)"),
        ("ChromaDB",      "Embedded RAG store",        "Pinecone (cloud)"),
        ("NetworkX",      "Python-native DAG",         "Neo4j (extra server)"),
        ("FastAPI + SSE", "Async streaming  ·  typed", "Flask (no async)"),
        ("Pydantic v2",   "Schema at boundaries",      "Raw dicts (no validation)"),
    ]
    widths = [Inches(2.6), Inches(5.0), Inches(4.5)]
    for j, w in enumerate(widths):
        tbl.columns[j].width = w
    for j, h in enumerate(headers):
        _style_table_header(tbl.cell(0, j), h)
    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            _style_table_body(tbl.cell(i, j), val, zebra=(i % 2 == 0), size=10)
    add_footer(slide, idx)


def build_pipeline(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Multi-Agent Pipeline", section="§4  AI Design")
    add_panel_title(slide, Inches(0.6), Inches(1.6), "SIX-STAGE PIPELINE", SKY, size=12)

    stages = [
        ("Parser",
         "BriefParser / TemplateParser  →  RequirementItem[]",
         SKY),
        ("Planner",
         "Rule decomposition + LLM reasoner → TaskList with type_reason · complexity_reason · tags",
         SKY),
        ("Effort",
         "RAG-calibrated (rule-dominant):  35 % RAG / 65 % rule  (C ≤ 2)  ·  15 % RAG / 85 % rule  (C ≥ 3)",
         AMBER),
        ("Critic",
         "Layer 1 rules + Layer 2 LLM  ·  Score < 0.50  ⟹  pipeline halts",
         ROSE),
        ("Graph",
         "NetworkX DAG · critical path · NFR-to-FR propagation (security · offline · compliance)",
         VIOLET),
        ("Risk + Sprint",
         "Six risk categories  ·  topological scheduling",
         EMERALD),
    ]
    top = Inches(2.0)
    h = Inches(0.7)
    for i, (head, body, color) in enumerate(stages):
        y = top + i * (h + Inches(0.05))
        add_card(slide, Inches(0.6), y, Inches(7.4), h)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.75), y + Inches(0.13), Inches(0.45), Inches(0.45))
        circle.line.fill.background()
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        ct = circle.text_frame
        ct.margin_left = ct.margin_right = 0
        ct.margin_top = ct.margin_bottom = 0
        ct.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ct.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i+1)
        cr.font.name = FONT_TITLE
        cr.font.size = Pt(14)
        cr.font.bold = True
        cr.font.color.rgb = BG_WHITE
        tb = slide.shapes.add_textbox(Inches(1.35), y + Inches(0.05), Inches(6.4), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = head
        r1.font.name = FONT_TITLE
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.font.color.rgb = INK_DEEP
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = FONT_BODY
        r2.font.size = Pt(10)
        r2.font.color.rgb = INK_BODY

    if PIPELINE_IMG.exists():
        add_card(slide, Inches(8.2), Inches(2.0), Inches(4.55), Inches(4.25),
                 fill=BG_WHITE, border=BORDER)
        slide.shapes.add_picture(str(PIPELINE_IMG),
                                  Inches(8.4), Inches(2.2), width=Inches(4.15))
        cap = slide.shapes.add_textbox(Inches(8.2), Inches(6.3), Inches(4.55), Inches(0.35))
        cp = cap.text_frame.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = "Pipeline architecture diagram"
        cr.font.name = FONT_BODY
        cr.font.size = Pt(10)
        cr.font.italic = True
        cr.font.color.rgb = INK_MUTED

    note = add_card(slide, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.55),
                    fill=SKY_LIGHT, border=SKY, border_w=0.5)
    add_text_in_card(note, [
        ("INTER-AGENT MODEL  ·  Agents communicate via typed Pydantic models (TaskList · CriticReport · RiskReport) — not raw JSON — validated at every boundary.",
         11, False, INK_DEEP, PP_ALIGN.CENTER, 0),
    ], padding_top=Inches(0.12))
    add_footer(slide, idx)


def build_hybrid(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Hybrid Reasoning: Rules + LLM + RAG", section="§4  AI Design")
    add_panel_title(slide, Inches(0.6), Inches(1.6), "THREE COOPERATING LAYERS", SKY, size=12)
    layers = [
        ("Rule Layer", "Deterministic Floor",
         "Regex extraction · keyword classification · complexity heuristics. Always runs first. Never fails.",
         EMERALD, "🛡️"),
        ("LLM Layer",  "Planning Reasoner",
         "Qwen 9B via Ollama. Receives the rule draft + retrieved KB context. Returns a refined plan.\n\nFallback: if Ollama unreachable or output fails schema validation, the rule plan ships unchanged.",
         SKY, "🤖"),
        ("RAG Layer",  "Knowledge Base · 100 docs",
         "Historical estimation (Desharnais · Maxwell) · COCOMO II calibration · Jones task-level patterns · domain planning examples.\n\nRetrieval: all-MiniLM-L6-v2 · top-2 by domain.",
         AMBER, "📚"),
    ]
    cw = Inches(3.93)
    ch = Inches(3.2)
    for i, (head, sub, body, color, icon) in enumerate(layers):
        left = Inches(0.6 + i * 4.07)
        add_card(slide, left, Inches(2.0), cw, ch)
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), cw, Inches(0.5))
        strip.line.fill.background()
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        ib = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.05), Inches(0.5), Inches(0.4))
        ip = ib.text_frame.paragraphs[0]
        ir = ip.add_run()
        ir.text = icon
        ir.font.name = "Segoe UI Emoji"
        ir.font.size = Pt(18)
        hb = slide.shapes.add_textbox(left + Inches(0.75), Inches(2.07), cw - Inches(0.85), Inches(0.4))
        hp = hb.text_frame.paragraphs[0]
        hr = hp.add_run()
        hr.text = head
        hr.font.name = FONT_TITLE
        hr.font.size = Pt(15)
        hr.font.bold = True
        hr.font.color.rgb = BG_WHITE
        sb = slide.shapes.add_textbox(left + Inches(0.25), Inches(2.65), cw - Inches(0.5), Inches(0.4))
        sp = sb.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = sub
        sr.font.name = FONT_BODY
        sr.font.size = Pt(12)
        sr.font.italic = True
        sr.font.color.rgb = INK_MUTED
        bb = slide.shapes.add_textbox(left + Inches(0.25), Inches(3.15),
                                       cw - Inches(0.5), ch - Inches(1.2))
        bf = bb.text_frame
        bf.word_wrap = True
        lines = body.split("\n")
        for li, line in enumerate(lines):
            bp = bf.paragraphs[0] if li == 0 else bf.add_paragraph()
            bp.line_spacing = 1.2
            br = bp.add_run()
            br.text = line
            br.font.name = FONT_BODY
            br.font.size = Pt(11)
            br.font.color.rgb = INK_BODY

    trace = add_card(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.5),
                     fill=SKY_LIGHT, border=SKY, border_w=1.0)
    add_text_in_card(trace, [
        ("🔍  TRACEABILITY GUARANTEE",   13, True, SKY_DARK, PP_ALIGN.LEFT, 0),
        ("Every run logs to plan_summary.json:   llm_attempted · llm_accepted · used_fallback · fallback_reason · retrieved_kb_context · kb_document_count",
         11, False, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("→  The AI's contribution is traceable, not hidden behind a black-box label.",
         11, False, INK_BODY, PP_ALIGN.LEFT, 0),
    ], padding_top=Inches(0.15))
    add_footer(slide, idx)


def build_evaluation(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Evaluation: Two Artifacts & Ablation", section="§5  Evaluation")

    # KPI tiles
    tile_w = Inches(2.93)
    tile_h = Inches(0.92)
    tiles_top = Inches(1.65)
    tiles = [
        ("100 %",  "Pass Rate (20 / 20)",     SKY),
        ("0.787",  "Overall System Score",    EMERALD),
        ("0.986",  "Avg Critic Score",        VIOLET),
        ("75 %",   "Direct LLM Acceptance",   AMBER),
    ]
    for i, (big, label, color) in enumerate(tiles):
        left = Inches(0.6 + i * 3.06)
        add_kpi_tile(slide, left, tiles_top, tile_w, tile_h, big, label, accent=color)

    # Left: Artifact details + ablation summary
    add_panel_title(slide, Inches(0.6), Inches(2.7), "ARTIFACTS  ·  ABLATION", SKY, size=12)
    text_card = add_card(slide, Inches(0.6), Inches(3.0), Inches(5.7), Inches(3.85))
    tf = text_card.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18)
    tf.word_wrap = True

    def _p(t, size, bold, color, italic=False, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        r = p.add_run()
        r.text = t
        r.font.name = FONT_BODY
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color

    _p("ARTIFACT A — Saved LLM + KB Run (v12)", 13, True, INK_DEEP, first=True)
    _p("MMRE / PRED(25)      0.359 / 0.668", 10, False, INK_BODY)
    _p("F1-FR / F1-NFR        0.705 / 0.470  ← NFR is acknowledged limit", 10, False, INK_BODY)
    _p("Coverage               1.000   ·   Fallback rate  25 %", 10, False, INK_BODY)
    _p(" ", 6, False, INK_BODY)
    _p("ARTIFACT B — Deterministic Rule Benchmark", 13, True, INK_DEEP)
    _p("FallbackOnlyClient forces rule-based mode  →", 10, False, INK_BODY)
    _p("byte-identical reruns without Ollama.", 10, False, INK_BODY)
    _p(" ", 6, False, INK_BODY)
    _p("Ablation — Honest Reporting", 13, True, AMBER)
    _p("R (Rules only)    Overall 0.749 · MMRE 0.271 · PRED25 0.778", 10, False, INK_BODY)
    _p("K (Rules + KB)    Overall 0.749 · MMRE 0.255 · PRED25 0.755", 10, False, INK_BODY)
    _p(" ", 6, False, INK_BODY)
    _p("↗  KB improves MMRE by 6 %  (better avg estimates)",  10, False, INK_BODY)
    _p("↘  KB lowers PRED(25) by 3 %  (slightly fewer ±25 %)", 10, False, INK_BODY)
    _p("Structural metrics unchanged · rule layer robust", 10, False, INK_BODY, italic=True)

    # Right: 2 charts
    add_threshold_chart(slide, Inches(6.45), Inches(2.7), Inches(6.3), Inches(2.1))
    add_ablation_chart(slide,  Inches(6.45), Inches(4.85), Inches(6.3), Inches(2.05))
    add_footer(slide, idx)


def build_closing(slide, idx):
    set_background(slide, BG_LIGHT)
    add_corner_accent(slide)
    add_title(slide, "Contributions · Limitations · Q & A", section="Closing")

    cw = Inches(3.93)
    ch = Inches(4.3)

    add_panel_title(slide, Inches(0.6), Inches(1.75), "RESEARCH CONTRIBUTIONS", SKY, size=12)
    card1 = add_card(slide, Inches(0.6), Inches(2.05), cw, ch, fill=BG_WHITE, border=SKY, border_w=1.0)
    add_text_in_card(card1, [
        ("①  A fully-local multi-agent",         11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("     project-planning prototype",      11, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("②  NFR-to-FR dependency-propagation",  11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("     rules across 4 NFR families",     11, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("③  Hybrid Rule + LLM + RAG",           11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("     with deterministic Critic",       11, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("     safety layer",                    11, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("④  Reproducible 20-sample benchmark",  11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("     + 2-condition ablation harness",  11, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("⑤  Auditable LLM reasoning trail",     11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("     in every output",                 11, False, INK_BODY, PP_ALIGN.LEFT, 0),
    ], padding_top=Inches(0.2))

    add_panel_title(slide, Inches(4.7), Inches(1.75), "BENEFICIARIES", EMERALD, size=12)
    card2 = add_card(slide, Inches(4.7), Inches(2.05), cw, ch, fill=BG_WHITE, border=EMERALD, border_w=1.0)
    add_text_in_card(card2, [
        ("•  Software-engineering students",     11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("    Structured planning training",     10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("•  Privacy-sensitive sectors",         11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("    Defense · healthcare · academia",  10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("•  Education researchers",             11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("    Open multi-agent baseline",        10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("•  Open-source community",             11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("    Reproducible benchmark",           10, False, INK_BODY, PP_ALIGN.LEFT, 0),
    ], padding_top=Inches(0.2))

    add_panel_title(slide, Inches(8.8), Inches(1.75), "ACKNOWLEDGED LIMITATIONS", ROSE, size=12)
    card3 = add_card(slide, Inches(8.8), Inches(2.05), cw, ch, fill=BG_WHITE, border=ROSE, border_w=1.0)
    add_text_in_card(card3, [
        ("⚠  F1-NFR  =  0.470",                   11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("    NFR boundary classification",       10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("    needs work",                        10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("⚠  No JWT-protected endpoints",         11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("    UI-layer role separation only",     10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("⚠  Monitor tested on synthetic Git",    11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("    Long-running repo pending",         10, False, INK_BODY, PP_ALIGN.LEFT, 0),
        (" ", 6, False, INK_BODY, PP_ALIGN.LEFT, 0),
        ("⚠  Single LLM evaluated (Qwen 9B)",     11, True, INK_DEEP, PP_ALIGN.LEFT, 0),
        ("    Broader LLM study future work",     10, False, INK_BODY, PP_ALIGN.LEFT, 0),
    ], padding_top=Inches(0.2))

    banner = add_card(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6),
                      fill=SKY_LIGHT, border=SKY, border_w=1.0)
    add_text_in_card(banner, [
        ("Thank you for your time.    I welcome your questions.",
         16, True, SKY_DARK, PP_ALIGN.CENTER, 0),
    ], padding_top=Inches(0.13))
    add_footer(slide, idx)


# ──────────────────────────────────────────────────────────────────────────────
# Speaker scripts — VERBATIM v2
# ──────────────────────────────────────────────────────────────────────────────
SCRIPTS = {
    1: (
        "Good morning. My graduation project is AI Project Manager — A Local Multi-Agent System for "
        "Automated Software Project Planning. Over fifteen minutes I will cover the problem, the "
        "multi-agent architecture, the hybrid reasoning approach, and the empirical evaluation on "
        "twenty project samples."
    ),
    2: (
        "Software projects fail not for lack of ideas but because converting a brief into an executable "
        "plan is manual, slow, and inconsistent. AI Project Manager closes that gap with four "
        "coordinated AI agents running locally on a Qwen nine-billion-parameter model. From a "
        "plain-text brief, the system produces: validated tasks with explicit reasoning fields, a "
        "dependency graph with critical-path analysis, calibrated effort estimates, a six-category "
        "risk report, and a sprint schedule. The architectural novelty is not the local LLM alone — "
        "it is the combination of multi-agent orchestration with a deterministic rule layer that "
        "makes every output auditable."
    ),
    3: "Five sections, then closing with contributions, limitations, and your questions.",
    4: (
        "The core idea is to replace the human project manager's first sprint — the planning sprint "
        "— with a coordinated AI team. Not a chatbot, not a template — a thinking system. It has "
        "seven capabilities. Each task carries explicit reasoning fields: type_reason explaining why "
        "it is functional or non-functional, and complexity_reason explaining its complexity score. "
        "Effort is calibrated through retrieval augmentation but keeps the rule-based estimate "
        "dominant. Dependencies are a real graph with critical-path computation. Risk is "
        "multi-category. And progress is monitored by matching real Git commits to tasks "
        "semantically."
    ),
    5: (
        "Four root causes motivate this work. First, knowledge fragmentation — every project's "
        "planning artifacts live in disconnected documents, with no machine-readable link between "
        "requirement, task, estimate, and risk. Second, subjective decomposition — give the same "
        "brief to two senior planners and you get two different plans, with no reproducible "
        "reasoning. Third, cloud lock-in — modern AI planning tools all require sending client "
        "requirements to external servers, which is unacceptable for regulated industries or "
        "sensitive academic projects. Fourth, black-box outputs — current tools produce plans "
        "without a reasoning trail, without a fallback when the model fails, and without an audit "
        "record that a stakeholder can defend. Our system addresses all four."
    ),
    6: (
        "Four strategic objectives shaped the design. First, build a comprehensive and auditable "
        "planner — every task carries its reasoning, the graph is validated, hours are calibrated, "
        "and risk is categorized. Second, compare AI architectures empirically — we measure the "
        "contribution of rules versus LLM versus RAG across twenty samples in an ablation study. "
        "Third, guarantee reliability through an independent Critic Agent with a hard rejection "
        "threshold. Fourth, optimize the trade-off between local performance and output quality — "
        "Mean Magnitude of Relative Error around zero-point-three-six on a single workstation, with "
        "no cloud dependency."
    ),
    7: (
        "The system rests on five pillars. Multi-agent orchestration — four agents, not one prompt, "
        "each with a narrow responsibility. Local-first operation — the Qwen model runs via Ollama "
        "on the same machine, with no paid API calls once the model weights and embedding files "
        "are cached. Hybrid reasoning — deterministic rules form the safety floor, the LLM acts as "
        "a planning reasoner, and a hundred-document ChromaDB knowledge base provides retrieval "
        "context. Dependency-aware output — we return a graph, not a list, with critical-path "
        "computation and a novel NFR-to-FR constraint propagation. And auditable evaluation — every "
        "claim is backed by twenty ground-truth samples and quantitative metrics."
    ),
    8: (
        "Why four agents and not one big prompt? A single LLM call is a competent writer but a "
        "poor project manager. It cannot simultaneously decompose tasks, validate its own "
        "consistency, propagate constraints, and score risk — without hallucinating in at least one "
        "of those dimensions. Our solution separates concerns. The Planner reads requirements and "
        "decomposes them with explicit reasoning fields. The Critic is an independent validator "
        "with rule-based and LLM layers — it has the authority to reject a plan and halt the "
        "pipeline if the score falls below zero-point-five. The Risk Agent scans six categories. "
        "The Monitor matches Git commits to tasks using sentence embeddings, with a keyword "
        "fallback. Each agent is independently testable and produces an auditable record."
    ),
    9: (
        "The literature traces three generations. The first — from nineteen-eighty-one — was "
        "algorithmic: COCOMO, Function Points. Mathematically defensible, but no task generation. "
        "The second applied classical machine learning to effort estimation and risk "
        "classification. Data-driven, but still no decomposition and no reasoning. The third "
        "generation, since twenty-twenty-three, brings large language models — GPT-Four-based "
        "tools, Jira AI, Linear Copilot. These finally accept natural-language input and produce "
        "task lists. But four open challenges remain: hallucination, privacy, no auditable "
        "reasoning, and no deterministic safety layer. Our system targets exactly those four."
    ),
    10: (
        "This table maps our contribution against the landscape. COCOMO and the PROMISE datasets "
        "do not decompose. Jira's AI, Linear's Copilot, and GitHub Workspace all use cloud LLMs, "
        "exposing client requirements to third parties and offering no auditable reasoning. "
        "AutoGPT-style agents lack determinism and validation. Our contribution closes four gaps: "
        "locality — a fully-local multi-agent planning prototype; auditability — explicit "
        "reasoning fields on every task; safety — an independent Critic with hard rejection "
        "authority; and NFR semantics — security, offline, and compliance non-functional "
        "requirements automatically propagate as constraints onto downstream functional tasks, a "
        "behavior absent from the surveyed literature."
    ),
    11: (
        "The system has two roles. The Student submits a brief and tracks progress: rating tasks "
        "one-to-five stars, marking status as todo, doing, or done, and asking the AI to explain "
        "any decision. The Supervisor reviews the plan — triaging tasks flagged for human review, "
        "approving or rejecting the full plan through an in-app chat, and exporting the committee "
        "brief as Excel. I want to be precise about scope here: the prototype uses a login page "
        "with demo credentials to separate roles at the UI layer through the Zustand store. "
        "Production deployment would require token-protected API endpoints — that is explicitly "
        "noted as future work."
    ),
    12: (
        "These tables summarize the requirements. On the functional side: parsing the brief, "
        "generating tasks via the multi-agent layer, validating through the Critic, calibrating "
        "effort against the knowledge base, building the dependency DAG, scheduling sprints, "
        "analyzing risk, and monitoring through Git. Eight functional requirements, each mapped "
        "to a specific module. On the non-functional side, six constraints: locality through no "
        "external API calls during operation, determinism through zero temperature and rule-based "
        "fallback, auditability through reasoning fields, reproducibility through a frozen "
        "twenty-sample benchmark, safety through the Critic's rejection threshold, and resilience "
        "through automatic context-window downscaling when the local model runs out of memory."
    ),
    13: (
        "The system follows a three-tier architecture. The presentation tier is a React-Nineteen "
        "single-page application with thirteen pages and live pipeline streaming via Server-Sent "
        "Events. The service tier is FastAPI on Python Three-Fourteen, fully async, fully typed, "
        "with ten routers exposing over thirty endpoints, all documented through Swagger. The "
        "intelligence tier is the local AI stack: Ollama serving Qwen at zero temperature for "
        "determinism, ChromaDB as the embedded vector store, MiniLM-L6 for embeddings, NetworkX "
        "for the dependency graph, and Pydantic-v2 to enforce schemas at every interface. Every "
        "technology choice has a documented rejection alternative — for example Ollama over "
        "GPT-Four for locality, NetworkX over Neo4j to avoid a second server."
    ),
    14: (
        "This is the six-stage pipeline. Stage one parses the brief into structured requirement "
        "items. Stage two is the Planner — it first applies rule-based decomposition, then if "
        "Ollama is available calls the LLM with the rule draft plus retrieved knowledge-base "
        "examples as a few-shot prompt. Every task carries explicit type-reason and "
        "complexity-reason fields. Stage three is effort estimation — and I want to be precise "
        "here: we blend retrieval-augmented estimates with rule-based estimates, but the rule "
        "estimate is always dominant. For simple tasks we use thirty-five percent RAG and "
        "sixty-five percent rule. For complex tasks we use only fifteen percent RAG and "
        "eighty-five percent rule, because complex tasks have noisier retrieval matches and we "
        "trust the deterministic rule more. Stage four is the Critic — two layers, rule-based "
        "first then LLM; if the score falls below zero-point-five the pipeline halts. Stage five "
        "builds the dependency graph and runs NFR-to-FR constraint propagation: a security NFR "
        "automatically blocks authentication and CRUD tasks; an offline-operation NFR blocks "
        "integration tasks. Stage six runs risk analysis and generates the sprint schedule. "
        "Agents communicate via typed Pydantic models, so schema drift is caught at every "
        "boundary."
    ),
    15: (
        "Hybrid reasoning is the architectural feature that distinguishes our system. Three "
        "layers cooperate. The rule layer is the deterministic floor — regex extraction, keyword "
        "classification, complexity heuristics. It always runs first and never fails. The LLM "
        "layer is the planning reasoner — Qwen-nine-billion served by Ollama, receiving the rule "
        "draft plus retrieved knowledge-base examples. If Ollama is unreachable or its output "
        "fails schema validation, the rule plan ships unchanged — graceful degradation, not "
        "crash. The RAG layer is a hundred-document ChromaDB knowledge base across four "
        "categories: historical estimation records from the Desharnais and Maxwell datasets, "
        "COCOMO calibration points, task-level patterns from Jones baselines, and "
        "domain-specific planning templates. Retrieval uses MiniLM embeddings, top-two by domain "
        "relevance. Every LLM call is logged with attempt status, acceptance status, fallback "
        "reason, and the exact retrieved context — making the AI contribution traceable, not "
        "hidden."
    ),
    16: (
        "Evaluation produced two distinct artifacts, and the distinction matters. Artifact A is "
        "a one-time recorded run using the full LLM-plus-knowledge-base configuration on twenty "
        "ground-truth samples across ten domains. It achieved a hundred-percent pass rate, an "
        "overall score of zero-point-seven-eight-seven, an average Critic score of "
        "zero-point-nine-eight-six, and a seventy-five-percent direct LLM acceptance rate — "
        "passing all four quality thresholds. I want to be honest about one number: the F1 score "
        "on non-functional classification is zero-point-four-seven-zero, which is an "
        "acknowledged limitation — NFR boundaries are inherently ambiguous and this is something "
        "we discuss in future work. Artifact B is the reproducible benchmark — the evaluation "
        "pipeline forces rule-based mode through a fallback-only client, so any reviewer can "
        "rerun it byte-identically without needing Ollama installed. On ablation: the knowledge "
        "base improves Mean Magnitude of Relative Error by six percent but actually lowers "
        "PRED-twenty-five by three percent. We report both directions honestly — the KB shifts "
        "estimates closer on average but slightly fewer fall within the strict twenty-five-"
        "percent band. Structural metrics are stable, which confirms the rule layer is robust "
        "under both conditions."
    ),
    17: (
        "To close, five research contributions: a fully-local multi-agent planning prototype, "
        "novel NFR-to-FR propagation rules across four NFR families, a hybrid architecture with "
        "a deterministic safety layer, a reproducible twenty-sample benchmark with ablation "
        "harness, and an auditable LLM reasoning trail. Beneficiaries include software-"
        "engineering students, privacy-sensitive sectors, education researchers, and the "
        "open-source community. I want to close honestly by stating four acknowledged "
        "limitations: NFR boundary classification needs work, backend endpoints are not yet "
        "JWT-protected, the Monitor Agent has been tested on synthetic data only, and we "
        "evaluated only one LLM model. These are clear directions for future work. Thank you — "
        "I welcome your questions."
    ),
}


BUILDERS = {
    1:  build_cover,
    2:  build_intro,
    3:  build_toc,
    4:  build_idea,
    5:  build_problem,
    6:  build_objectives,
    7:  build_pillars,
    8:  build_multiagent,
    9:  build_generations,
    10: build_relwork,
    11: build_actors,
    12: build_requirements,
    13: build_architecture,
    14: build_pipeline,
    15: build_hybrid,
    16: build_evaluation,
    17: build_closing,
}


def load_real_numbers():
    v12 = json.loads((REPO / "data/evaluation/evaluation_llm_kb_report_v12.json")
                     .read_text(encoding="utf-8"))
    ab = json.loads((REPO / "data/evaluation/ablation_report.json")
                    .read_text(encoding="utf-8"))
    return {
        "pass_rate": v12["pass_rate_pct"],
        "overall":   v12["metric_summary"]["average_overall_score"],
        "critic":    v12["average_critic_score"],
        "direct":    v12["direct_generation_rate_pct"],
        "mmre":      v12["metric_summary"]["average_mmre"],
        "pred25":    v12["metric_summary"]["average_pred25"],
        "f1_fr":     v12["metric_summary"]["average_f1_fr"],
        "f1_nfr":    v12["metric_summary"]["average_f1_nfr"],
        "R_mmre":    ab["conditions"]["R"]["average_mmre"],
        "K_mmre":    ab["conditions"]["K"]["average_mmre"],
        "R_pred25":  ab["conditions"]["R"]["average_pred25"],
        "K_pred25":  ab["conditions"]["K"]["average_pred25"],
    }


def main():
    nums = load_real_numbers()
    print("-" * 70)
    print("Evaluation numbers locked in:")
    for k, v in nums.items():
        print(f"  {k:12s} = {v}")
    print("-" * 70)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for i in range(1, 18):
        slide = prs.slides.add_slide(blank)
        builder = BUILDERS[i]
        if i == 1:
            builder(slide)
        else:
            builder(slide, i)
        add_notes(slide, SCRIPTS[i])
        add_slide_transition(slide, kind="fade")
        # NO cascade animation — academic version
        print(f"  [{i:2d}/17] {builder.__name__:22s}  notes={len(SCRIPTS[i]):4d}c")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print("-" * 70)
    print(f"Saved -> {OUT_PATH}")
    print(f"Slides: 17  ·  v4 visual design  ·  v2 verbatim text  ·  NO animations")
    print(f"File size: {OUT_PATH.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
